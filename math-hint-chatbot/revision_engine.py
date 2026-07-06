import os
import re
import fitz  # PyMuPDF
import docx
import pptx

class RevisionEngine:

    def __init__(self):

        self.upload_folder = "uploads"

        os.makedirs(self.upload_folder, exist_ok=True)

    # -----------------------------------
    # Read TXT
    # -----------------------------------

    def read_txt(self, filepath):

        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:

            return f.read()

    # -----------------------------------
    # Read PDF
    # -----------------------------------

    def read_pdf(self, filepath):

        text = ""

        pdf = fitz.open(filepath)

        for page in pdf:

            text += page.get_text()

        pdf.close()

        return text

    # -----------------------------------
    # Read DOCX
    # -----------------------------------

    def read_docx(self, filepath):

        document = docx.Document(filepath)

        text = []

        for para in document.paragraphs:

            text.append(para.text)

        return "\n".join(text)

    # -----------------------------------
    # Read PPTX
    # -----------------------------------

    def read_pptx(self, filepath):

        presentation = pptx.Presentation(filepath)

        text = []

        for slide in presentation.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    text.append(shape.text)

        return "\n".join(text)

    # -----------------------------------
    # Extract File Text
    # -----------------------------------

    def extract_text(self, filepath):

        extension = os.path.splitext(filepath)[1].lower()

        if extension == ".pdf":

            return self.read_pdf(filepath)

        elif extension == ".docx":

            return self.read_docx(filepath)

        elif extension == ".pptx":

            return self.read_pptx(filepath)

        elif extension == ".txt":

            return self.read_txt(filepath)

        else:

            raise Exception("Unsupported file format.")

    # -----------------------------------
    # Clean Text
    # -----------------------------------

    def clean_text(self, text):

        text = re.sub(r"\s+", " ", text)

        return text.strip()
        # -----------------------------------
    # Generate Summary
    # -----------------------------------

    def generate_summary(self, text):

        sentences = re.split(r'(?<=[.!?]) +', text)

        summary = " ".join(sentences[:8])

        if not summary.strip():

            summary = text[:1000]

        return summary

    # -----------------------------------
    # Extract Formula-like Expressions
    # -----------------------------------

    def extract_formulas(self, text):

        pattern = r"[A-Za-z0-9\^\+\-\*/=\(\)]+"

        matches = re.findall(pattern, text)

        formulas = []

        for item in matches:

            if any(op in item for op in ["=", "+", "-", "*", "/", "^"]):

                if len(item) > 2:

                    formulas.append(item)

        formulas = list(dict.fromkeys(formulas))

        return formulas[:20]

    # -----------------------------------
    # Extract Important Topics
    # -----------------------------------

    def extract_topics(self, text):

        words = re.findall(r"[A-Za-z]+", text)

        ignore = {

            "the","is","and","of","to","a","an","for","on",

            "in","with","by","that","this","are","be",

            "from","or","as","at","it","we","you"

        }

        frequency = {}

        for word in words:

            word = word.lower()

            if len(word) < 4:

                continue

            if word in ignore:

                continue

            frequency[word] = frequency.get(word, 0) + 1

        topics = sorted(

            frequency.items(),

            key=lambda x: x[1],

            reverse=True

        )

        return [topic[0].title() for topic in topics[:10]]

    # -----------------------------------
    # Extract Keywords
    # -----------------------------------

    def extract_keywords(self, text):

        topics = self.extract_topics(text)

        return topics[:15]

    # -----------------------------------
    # Estimate Reading Time
    # -----------------------------------

    def reading_time(self, text):

        words = len(text.split())

        minutes = max(1, round(words / 180))

        return f"{minutes} min"
        # -----------------------------------
    # Generate Flashcards
    # -----------------------------------

    def generate_flashcards(self, topics):

        flashcards = []

        for topic in topics[:10]:

            flashcards.append({

                "question": f"What is {topic}?",

                "answer": f"Revise the concept of {topic} from your notes."

            })

        return flashcards

    # -----------------------------------
    # Generate MCQs
    # -----------------------------------

    def generate_mcqs(self, topics):

        mcqs = []

        for topic in topics[:5]:

            mcqs.append({

                "question": f"Which topic best describes '{topic}'?",

                "options": [

                    topic,

                    "Probability",

                    "Integration",

                    "Matrices"

                ],

                "answer": topic

            })

        return mcqs

    # -----------------------------------
    # Identify Weak Topics
    # -----------------------------------

    def identify_weak_topics(self, topics):

        if len(topics) <= 3:

            return topics

        return topics[-3:]

    # -----------------------------------
    # Generate Revision Plan
    # -----------------------------------

    def generate_revision_plan(self, topics):

        plan = []

        day = 1

        for topic in topics[:7]:

            plan.append(

                f"Day {day}: Revise {topic}"

            )

            day += 1

        return plan

    # -----------------------------------
    # Analyze Notes
    # -----------------------------------

    def analyze(self, filepath):

        text = self.extract_text(filepath)

        text = self.clean_text(text)

        summary = self.generate_summary(text)

        topics = self.extract_topics(text)

        formulas = self.extract_formulas(text)

        flashcards = self.generate_flashcards(topics)

        mcqs = self.generate_mcqs(topics)

        weak_topics = self.identify_weak_topics(topics)

        revision_plan = self.generate_revision_plan(topics)

        return {

            "status": "success",

            "summary": summary,

            "topics": topics,

            "keywords": self.extract_keywords(text),

            "formulas": formulas,

            "reading_time": self.reading_time(text),

            "flashcards": flashcards,

            "mcqs": mcqs,

            "weak_topics": weak_topics,

            "revision_plan": revision_plan

        }
